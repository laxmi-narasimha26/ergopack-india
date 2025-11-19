#!/usr/bin/env python3
"""
ErgoPack India - Premium Pitch Deck Generator
Creates a professional investor pitch deck with advanced design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ErgoPack Brand Colors
ERGOPACK_RED = RGBColor(200, 16, 46)      # #C8102E
ERGOPACK_GOLD = RGBColor(255, 184, 28)    # #FFB81C
CHARCOAL = RGBColor(45, 45, 45)           # #2D2D2D
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_GRAY = RGBColor(80, 80, 80)

def create_presentation():
    """Create the premium pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title/Cover
    create_title_slide(prs)

    # Slide 2: The Problem
    create_problem_slide(prs)

    # Slide 3: The Solution
    create_solution_slide(prs)

    # Slide 4: Why ErgoPack
    create_why_ergopack_slide(prs)

    # Slide 5: 27-Year Heritage
    create_heritage_slide(prs)

    # Slide 6: Timeline & Milestones
    create_timeline_slide(prs)

    # Slide 7: Product Portfolio
    create_product_portfolio_slide(prs)

    # Slide 8: X-pert Line (Premium)
    create_xpert_line_slide(prs)

    # Slide 9: Technology Innovation - ChainLance
    create_chainlance_slide(prs)

    # Slide 10: Lithium-Ion Advantage
    create_lithium_slide(prs)

    # Slide 11: Target Markets
    create_target_markets_slide(prs)

    # Slide 12: Industry Solutions - Pharma
    create_pharma_slide(prs)

    # Slide 13: Global Traction
    create_traction_slide(prs)

    # Slide 14: India Partnership
    create_india_partnership_slide(prs)

    # Slide 15: Competitive Advantages
    create_competitive_advantages_slide(prs)

    # Slide 16: Vision & Future
    create_vision_slide(prs)

    # Slide 17: Call to Action
    create_cta_slide(prs)

    return prs

def add_background(slide, color=CHARCOAL):
    """Add colored background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, top=0.5, left=0.5, width=9, height=1.5,
                   color=WHITE, size=54, bold=True, align=PP_ALIGN.LEFT):
    """Add title text box"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return textbox

def add_body_text(slide, text, top, left=0.5, width=9, height=1,
                  color=WHITE, size=18, align=PP_ALIGN.LEFT):
    """Add body text box"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = align
    return textbox

def add_bullet_points(slide, points, top, left=0.8, width=8.5, color=WHITE, size=16):
    """Add bullet point list"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(12)

    return textbox

def add_accent_bar(slide, top, color=ERGOPACK_RED):
    """Add decorative accent bar"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(top),
        Inches(0.15), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def create_title_slide(prs):
    """Slide 1: Title/Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)

    # Large Title
    add_title_text(slide, "ERGOPACK", top=2, left=0.5, size=72, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    # Subtitle
    add_title_text(slide, "INDIA", top=2.8, left=0.5, size=64, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)

    # Tagline
    add_body_text(slide, "German Precision Engineering for Cargo Security",
                  top=4.2, left=0.5, width=9, size=22, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER)

    # Footer
    add_body_text(slide, "27 Years of Innovation | 15,000+ Systems | 60+ Countries",
                  top=6.5, left=0.5, width=9, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def create_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_accent_bar(slide, 0.5)
    add_title_text(slide, "The Problem", top=0.5, size=48, color=ERGOPACK_RED)

    add_body_text(slide, "High-Consequence Industries Face Critical Cargo Security Risks",
                  top=1.8, size=24, color=WHITE)

    points = [
        "📦 Pharmaceutical shipments: Temperature-sensitive cargo worth millions at risk",
        "🚗 Automotive JIT delivery: One delayed shipment = entire assembly line stoppage",
        "💻 Electronics: High-value fragile goods demanding zero damage tolerance",
        "⚠️ Load shifting during transit causes $2.8B+ annual losses in India",
        "📊 Manual strapping inconsistency creates unverifiable load integrity",
        "👷 Traditional methods are slow, unreliable, and damage-prone"
    ]

    add_bullet_points(slide, points, top=2.8, size=18)

def create_solution_slide(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_accent_bar(slide, 0.5, ERGOPACK_GOLD)
    add_title_text(slide, "The Solution", top=0.5, size=48, color=ERGOPACK_GOLD)

    add_body_text(slide, "Premium Automated Strapping Systems for Mission-Critical Cargo",
                  top=1.8, size=24, color=WHITE)

    points = [
        "🎯 Precision Automated Strapping: Consistent tension, perfect placement, every time",
        "🔬 German-Engineered Quality: 27 years of continuous innovation from ErgoPack Germany",
        "⚡ Patented ChainLance Technology: World's first precision strap guidance system",
        "🔋 Advanced Lithium-Ion Power: 1,200 cycles per charge (3.4x more than competitors)",
        "📱 Siemens Touchscreen Controls: Industry 4.0 ready with IoT integration",
        "✅ Verifiable Load Integrity: Automatic documentation for compliance & audits"
    ]

    add_bullet_points(slide, points, top=2.8, size=18)

def create_why_ergopack_slide(prs):
    """Slide 4: Why ErgoPack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_accent_bar(slide, 0.5, ERGOPACK_RED)
    add_title_text(slide, "Why ErgoPack?", top=0.5, size=48, color=ERGOPACK_RED)

    add_body_text(slide, "Not a Tool. A Strategic Risk Mitigation System.",
                  top=1.8, size=28, color=ERGOPACK_GOLD)

    points = [
        "🏆 Award-Winning Innovation: Gold Medal Geneva Exhibition + German Inventor Award",
        "🔐 Patented Technology: ChainLance system (1999) - world's first, still unmatched",
        "🌍 Global Market Leader: 15,000+ systems operating in 60+ countries",
        "📈 99.8% Customer Satisfaction: Trusted by Fortune 500 companies worldwide",
        "🏭 Made in Germany: Uncompromising precision, reliability, and performance",
        "🎓 ISO Certified: ISO 12100:2010, EU Conformity, AGR Ergonomic Certification"
    ]

    add_bullet_points(slide, points, top=2.8, size=18)

def create_heritage_slide(prs):
    """Slide 5: 27-Year Heritage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_title_text(slide, "27 Years of German Innovation", top=0.5, size=48, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)

    add_body_text(slide, "From Prototype to Global Leader",
                  top=1.5, size=22, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    # Key milestones in a visually appealing format
    milestones = [
        ("1998", "Revolutionary double scissors lifting mechanism prototype"),
        ("1999", "ChainLance invented & patented - world's first"),
        ("2002", "🏅 Gold Medal - Geneva International Exhibition of Inventions"),
        ("2008", "First integrated sealing systems (720E/730E)"),
        ("2014", "ErgoPack Air - world's only mobile elevated strapping system"),
        ("2018", "Touchscreen innovation - Siemens industrial interface"),
        ("2025", "Global expansion: 15,000+ systems in 60+ countries")
    ]

    top_pos = 2.5
    for year, desc in milestones:
        # Year
        year_box = slide.shapes.add_textbox(Inches(1), Inches(top_pos), Inches(1.2), Inches(0.4))
        p = year_box.text_frame.paragraphs[0]
        p.text = year
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ERGOPACK_GOLD

        # Description
        desc_box = slide.shapes.add_textbox(Inches(2.5), Inches(top_pos), Inches(7), Inches(0.4))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        top_pos += 0.6

def create_timeline_slide(prs):
    """Slide 6: Timeline & Milestones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_title_text(slide, "Innovation Timeline", top=0.5, size=48, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    timeline_events = [
        ("1998-2001", "Foundation Era", "First prototypes, ChainLance patent, first production models"),
        ("2002-2008", "Recognition Era", "International awards, electronic revolution, Tool-Lift system"),
        ("2011-2014", "Advanced Tech Era", "State-of-art sealing, ErgoPack Air mobile system"),
        ("2016-2019", "Digital Era", "40+ new features, touchscreens, X-pert Line benchmark"),
        ("2020-2025", "Global Expansion", "15,000+ systems, 60+ countries, India partnership")
    ]

    top_pos = 2
    for period, era, desc in timeline_events:
        # Period
        period_box = slide.shapes.add_textbox(Inches(0.8), Inches(top_pos), Inches(1.8), Inches(0.35))
        p = period_box.text_frame.paragraphs[0]
        p.text = period
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ERGOPACK_RED

        # Era
        era_box = slide.shapes.add_textbox(Inches(2.8), Inches(top_pos), Inches(2.5), Inches(0.35))
        p = era_box.text_frame.paragraphs[0]
        p.text = era
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ERGOPACK_GOLD

        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.5), Inches(top_pos), Inches(4), Inches(0.35))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

        top_pos += 1

def create_product_portfolio_slide(prs):
    """Slide 7: Product Portfolio"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_title_text(slide, "Complete Product Portfolio", top=0.5, size=48, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)
    add_body_text(slide, "11 Models Across 4 Product Lines", top=1.3, size=24, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # Product lines
    lines = [
        {
            "name": "X-PERT LINE (Premium)",
            "models": "4 Models",
            "features": "Lithium-Ion • Siemens Touchscreen • Line Laser • 1,200 cycles/charge",
            "color": ERGOPACK_GOLD
        },
        {
            "name": "ECONOMY LINE (Reliable)",
            "models": "5 Models",
            "features": "Lead-Fleece Battery • Proven Technology • 350 cycles/charge",
            "color": WHITE
        },
        {
            "name": "MOBILE SYSTEMS (Flexible)",
            "models": "2 Models",
            "features": "Portable • Retracting • Elevated Strapping • Versatile Mounting",
            "color": ERGOPACK_RED
        }
    ]

    top_pos = 2.5
    for line in lines:
        # Line name
        name_box = slide.shapes.add_textbox(Inches(1), Inches(top_pos), Inches(8), Inches(0.4))
        p = name_box.text_frame.paragraphs[0]
        p.text = line["name"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = line["color"]

        # Models count
        models_box = slide.shapes.add_textbox(Inches(1.5), Inches(top_pos + 0.4), Inches(2), Inches(0.3))
        p = models_box.text_frame.paragraphs[0]
        p.text = line["models"]
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY

        # Features
        features_box = slide.shapes.add_textbox(Inches(1.5), Inches(top_pos + 0.7), Inches(7), Inches(0.5))
        tf = features_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = line["features"]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        top_pos += 1.5

def create_xpert_line_slide(prs):
    """Slide 8: X-pert Line Premium"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_accent_bar(slide, 0.5, ERGOPACK_GOLD)
    add_title_text(slide, "X-PERT LINE", top=0.5, size=48, color=ERGOPACK_GOLD)
    add_body_text(slide, "The Premium Professional Series", top=1.3, size=20, color=WHITE)

    points = [
        "🔋 Lithium-Ion Power: 1,200 strapping cycles per charge (industry-leading)",
        "📱 Siemens Industrial Touchscreen: Intuitive control & real-time monitoring",
        "🎯 Line Laser Type 2: Precision positioning for perfect strap placement",
        "⚡ 66 m/min Chain Speed: Ultra-fast operation (65% faster than economy)",
        "🎚️ Adjustable Tension: 150N - 4,500N range (covers all applications)",
        "⚙️ 4 Specialized Models: 745X Li, 726X Li, 713X Li, 700X Li",
        "🌐 IoT-Ready: Industry 4.0 integration for smart factory environments"
    ]

    add_bullet_points(slide, points, top=2.2, size=17)

def create_chainlance_slide(prs):
    """Slide 9: ChainLance Technology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_title_text(slide, "CHAINLANCE TECHNOLOGY", top=0.5, size=48, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)
    add_body_text(slide, "Patented 1999 • World's First • Still Unmatched", top=1.3, size=20, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    add_body_text(slide, "The Core Innovation That Revolutionized Automated Strapping",
                  top=2, size=24, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    points = [
        "✓ Precision Strap Guidance: Guides strap around pallet with millimeter accuracy",
        "✓ Zero Load Shifting: Eliminates vibration and movement during transit",
        "✓ Perfect Placement: Consistent strap position every single cycle",
        "✓ Multi-Material Compatible: PP, PET, Paper, Cord, Composite straps",
        "✓ Reduced Damage: Prevents cargo shifting that causes $2.8B+ annual losses",
        "✓ Verifiable Results: Consistent, documentable load securing for compliance"
    ]

    add_bullet_points(slide, points, top=3, size=18)

    add_body_text(slide, "The ChainLance system is the hallmark of every ErgoPack system worldwide",
                  top=6.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, width=9)

def create_lithium_slide(prs):
    """Slide 10: Lithium-Ion Advantage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_title_text(slide, "Lithium-Ion Advantage", top=0.5, size=48, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # Comparison boxes
    # Lead-Fleece (Old)
    old_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2),
        Inches(4), Inches(4)
    )
    old_box.fill.solid()
    old_box.fill.fore_color.rgb = DARK_GRAY
    old_box.line.color.rgb = WHITE

    old_title = slide.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(3), Inches(0.4))
    p = old_title.text_frame.paragraphs[0]
    p.text = "Lead-Fleece Battery"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    old_specs = [
        "350 cycles/charge",
        "8-10 hours charging",
        "12.3 kg weight",
        "Legacy technology"
    ]

    spec_top = 3
    for spec in old_specs:
        spec_box = slide.shapes.add_textbox(Inches(1.3), Inches(spec_top), Inches(3), Inches(0.3))
        p = spec_box.text_frame.paragraphs[0]
        p.text = spec
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        spec_top += 0.5

    # Lithium-Ion (New)
    new_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(2),
        Inches(4), Inches(4)
    )
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = ERGOPACK_RED
    new_box.line.color.rgb = ERGOPACK_GOLD
    new_box.line.width = Pt(3)

    new_title = slide.shapes.add_textbox(Inches(5.7), Inches(2.3), Inches(3), Inches(0.4))
    p = new_title.text_frame.paragraphs[0]
    p.text = "Lithium-Ion Battery"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ERGOPACK_GOLD
    p.alignment = PP_ALIGN.CENTER

    new_specs = [
        "1,200 cycles/charge ⚡",
        "3.5 hours charging ⏱️",
        "5 kg weight 🪶",
        "Advanced technology 🚀"
    ]

    spec_top = 3
    for spec in new_specs:
        spec_box = slide.shapes.add_textbox(Inches(5.7), Inches(spec_top), Inches(3), Inches(0.3))
        p = spec_box.text_frame.paragraphs[0]
        p.text = spec
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        spec_top += 0.5

    # Bottom comparison
    add_body_text(slide, "3.4x More Cycles • 60% Faster Charging • 60% Lighter",
                  top=6.5, size=20, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

def create_target_markets_slide(prs):
    """Slide 11: Target Markets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_accent_bar(slide, 0.5, ERGOPACK_RED)
    add_title_text(slide, "Target Markets: The Elite 60", top=0.5, size=48, color=ERGOPACK_RED)
    add_body_text(slide, "High-Consequence Industries Where Failure is Not an Option", top=1.5, size=20, color=WHITE)

    markets = [
        {
            "icon": "💊",
            "name": "Pharmaceuticals",
            "value": "$42B market in India",
            "needs": "GDP/GMP compliance • Temperature-sensitive • Zero damage tolerance"
        },
        {
            "icon": "🚗",
            "name": "Automotive Components",
            "value": "$67B market in India",
            "needs": "JIT delivery critical • Assembly line dependency • Quality standards"
        },
        {
            "icon": "💻",
            "name": "Electronics & Machinery",
            "value": "$118B market in India",
            "needs": "High-value goods • Fragile components • Export quality required"
        }
    ]

    top_pos = 2.5
    for market in markets:
        # Icon and name
        icon_box = slide.shapes.add_textbox(Inches(1), Inches(top_pos), Inches(0.5), Inches(0.4))
        p = icon_box.text_frame.paragraphs[0]
        p.text = market["icon"]
        p.font.size = Pt(32)

        name_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_pos), Inches(4), Inches(0.4))
        p = name_box.text_frame.paragraphs[0]
        p.text = market["name"]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ERGOPACK_GOLD

        # Value
        value_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_pos + 0.4), Inches(4), Inches(0.3))
        p = value_box.text_frame.paragraphs[0]
        p.text = market["value"]
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

        # Needs
        needs_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_pos + 0.7), Inches(7.5), Inches(0.5))
        tf = needs_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = market["needs"]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        top_pos += 1.6

def create_pharma_slide(prs):
    """Slide 12: Pharma Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_title_text(slide, "Pharmaceutical Solution", top=0.5, size=48, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)

    add_body_text(slide, "Case Study: Temperature-Sensitive Cargo Protection",
                  top=1.4, size=22, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # The Challenge
    challenge_title = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4), Inches(0.4))
    p = challenge_title.text_frame.paragraphs[0]
    p.text = "THE CHALLENGE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ERGOPACK_RED

    challenges = [
        "• $5M+ shipment value at risk",
        "• GDP/GMP compliance required",
        "• Temperature excursions = total loss",
        "• Manual strapping inconsistent",
        "• No audit trail for compliance"
    ]

    challenge_top = 2.7
    for challenge in challenges:
        cb = slide.shapes.add_textbox(Inches(1), Inches(challenge_top), Inches(3.5), Inches(0.25))
        p = cb.text_frame.paragraphs[0]
        p.text = challenge
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        challenge_top += 0.35

    # The Solution
    solution_title = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4), Inches(0.4))
    p = solution_title.text_frame.paragraphs[0]
    p.text = "THE ERGOPACK SOLUTION"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ERGOPACK_GOLD

    solutions = [
        "✓ ChainLance precision: zero shifting",
        "✓ Automatic compliance docs",
        "✓ 99.99% reliability rate",
        "✓ Real-time load monitoring",
        "✓ Verifiable integrity for audits"
    ]

    solution_top = 2.7
    for solution in solutions:
        sb = slide.shapes.add_textbox(Inches(5.4), Inches(solution_top), Inches(3.5), Inches(0.25))
        p = sb.text_frame.paragraphs[0]
        p.text = solution
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        solution_top += 0.35

    # Results
    add_body_text(slide, "RESULTS: Zero load failures • 100% compliance • $2.3M saved annually",
                  top=5.5, size=18, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

def create_traction_slide(prs):
    """Slide 13: Global Traction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_title_text(slide, "Global Traction & Proof Points", top=0.5, size=48, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # Big numbers
    metrics = [
        {"number": "15,000+", "label": "Systems Deployed", "color": ERGOPACK_RED},
        {"number": "60+", "label": "Countries", "color": ERGOPACK_GOLD},
        {"number": "99.8%", "label": "Customer Satisfaction", "color": WHITE},
        {"number": "27", "label": "Years Innovation", "color": ERGOPACK_RED}
    ]

    left_pos = 0.8
    for metric in metrics:
        # Number
        num_box = slide.shapes.add_textbox(Inches(left_pos), Inches(2), Inches(2), Inches(0.8))
        p = num_box.text_frame.paragraphs[0]
        p.text = metric["number"]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = metric["color"]
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(left_pos), Inches(2.8), Inches(2), Inches(0.4))
        p = label_box.text_frame.paragraphs[0]
        p.text = metric["label"]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        left_pos += 2.3

    # Additional proof points
    add_body_text(slide, "Trusted by Fortune 500 Companies Worldwide",
                  top=4, size=22, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    points = [
        "🏆 Gold Medal - Geneva International Exhibition of Inventions",
        "🥇 German Inventor's Award Winner",
        "✓ ISO 12100:2010 Machine Safety Certified",
        "✓ AGR Certified for Ergonomic Excellence",
        "✓ EU Declaration of Conformity"
    ]

    add_bullet_points(slide, points, top=4.7, left=1.5, width=7, size=16)

def create_india_partnership_slide(prs):
    """Slide 14: India Partnership"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_accent_bar(slide, 0.5, ERGOPACK_GOLD)
    add_title_text(slide, "Strategic India Partnership", top=0.5, size=48, color=ERGOPACK_GOLD)

    add_body_text(slide, "Bringing German Precision to India's High-Growth Industries",
                  top=1.5, size=22, color=WHITE)

    # Partnership details
    add_body_text(slide, "🤝 BENZ PACKAGING SOLUTIONS", top=2.3, size=24, color=ERGOPACK_RED)

    points = [
        "✓ Exclusive distribution partner for ErgoPack in India",
        "✓ Local service & support infrastructure across major industrial hubs",
        "✓ Custom solutions for Indian pharmaceutical, automotive & electronics sectors",
        "✓ Training & certification programs for technical teams",
        "✓ Spare parts inventory & 24/7 customer support"
    ]

    add_bullet_points(slide, points, top=3, size=17)

    # Market opportunity
    add_body_text(slide, "INDIA MARKET OPPORTUNITY", top=5, size=20, color=ERGOPACK_GOLD)

    market_points = [
        "• $227B+ addressable market (Pharma + Auto + Electronics)",
        "• Fastest-growing manufacturing hub in Asia-Pacific",
        "• Government 'Make in India' initiative driving quality standards",
        "• Rising demand for GDP/GMP compliance in pharmaceutical sector"
    ]

    add_bullet_points(slide, market_points, top=5.5, size=15)

def create_competitive_advantages_slide(prs):
    """Slide 15: Competitive Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CHARCOAL)

    add_title_text(slide, "Competitive Advantages", top=0.5, size=48, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)
    add_body_text(slide, "Why We Win Against Competitors", top=1.3, size=20, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    advantages = [
        "🔐 Patented ChainLance: Proprietary technology competitors cannot replicate",
        "🏭 Made in Germany: Superior build quality & engineering precision",
        "🔋 Lithium-Ion Leadership: 3.4x more cycles than competitor systems",
        "📱 Industry 4.0 Ready: Siemens touchscreen & IoT integration",
        "🏆 Award-Winning: International recognition for innovation excellence",
        "🌍 Global Scale: 15,000+ installations provide extensive field-proven reliability",
        "✅ Complete Portfolio: 11 models covering every application need",
        "🎯 Premium Positioning: Not a commodity tool - strategic risk mitigation",
        "📊 Verifiable Results: Automatic documentation for compliance",
        "🛠️ Total Solution: Equipment + service + training + support ecosystem"
    ]

    add_bullet_points(slide, advantages, top=2.2, size=16)

def create_vision_slide(prs):
    """Slide 16: Vision & Future"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BLACK)

    add_title_text(slide, "Vision for India", top=1.5, size=54, color=ERGOPACK_RED, align=PP_ALIGN.CENTER, width=9)

    add_body_text(slide, "Becoming India's Trusted Partner for Mission-Critical Cargo Security",
                  top=2.5, size=24, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # Vision pillars
    pillars = [
        ("Market Leadership", "Capture 40% of premium automated strapping market by 2027"),
        ("Innovation Hub", "Establish India R&D center for Asia-Pacific innovations"),
        ("Customer Success", "500+ installations across pharma, auto & electronics by 2026"),
        ("Service Excellence", "Build nationwide service network with <4 hour response time")
    ]

    top_pos = 3.8
    for title, desc in pillars:
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(top_pos), Inches(7), Inches(0.3))
        p = title_box.text_frame.paragraphs[0]
        p.text = f"• {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ERGOPACK_GOLD

        # Description
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(top_pos + 0.35), Inches(6.5), Inches(0.3))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        top_pos += 0.8

def create_cta_slide(prs):
    """Slide 17: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, ERGOPACK_RED)

    add_title_text(slide, "Partner With Excellence", top=2, size=54, color=WHITE, align=PP_ALIGN.CENTER, width=9)

    add_body_text(slide, "Secure Your Cargo. Secure Your Future.",
                  top=3, size=32, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

    # Contact info
    contact_lines = [
        "📧 info@ergopack-india.com",
        "🌐 www.ergopack-india.com",
        "📞 +91 XXX XXX XXXX",
        "",
        "🤝 Strategic Partnership: Benz Packaging Solutions"
    ]

    top_pos = 4.5
    for line in contact_lines:
        cb = slide.shapes.add_textbox(Inches(2), Inches(top_pos), Inches(6), Inches(0.3))
        p = cb.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        top_pos += 0.4

    add_body_text(slide, "German Engineering • Indian Innovation • Global Excellence",
                  top=6.8, size=16, color=ERGOPACK_GOLD, align=PP_ALIGN.CENTER, width=9)

def main():
    """Main execution"""
    print("Creating ErgoPack India Premium Pitch Deck...")

    prs = create_presentation()

    output_file = "/home/user/ergopack-india/ErgoPack_India_Pitch_Deck.pptx"
    prs.save(output_file)

    print(f"✓ Pitch deck created successfully!")
    print(f"✓ File saved: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("\nSlide Breakdown:")
    print("  1. Title/Cover")
    print("  2. The Problem")
    print("  3. The Solution")
    print("  4. Why ErgoPack")
    print("  5. 27-Year Heritage")
    print("  6. Timeline & Milestones")
    print("  7. Product Portfolio")
    print("  8. X-pert Line (Premium)")
    print("  9. ChainLance Technology")
    print("  10. Lithium-Ion Advantage")
    print("  11. Target Markets")
    print("  12. Industry Solutions - Pharma")
    print("  13. Global Traction")
    print("  14. India Partnership")
    print("  15. Competitive Advantages")
    print("  16. Vision & Future")
    print("  17. Call to Action")

if __name__ == "__main__":
    main()
